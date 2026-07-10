"""
attachment_service.py — Shared file context pipeline for multi-teammate collaboration.

Pipeline:
  File bytes → Parse → AttachmentContext → Cache → Inject into ALL teammate prompts

All teammates receive the SAME parsed context snapshot (no individual parsing).
"""
import logging
import os
import re
import hashlib
from typing import Any, Optional
from dataclasses import dataclass, field, asdict

from backend.cache import TTLCache

logger = logging.getLogger("attachment_service")

# ── Shared cache (file_id → AttachmentContext dict) ──
attachment_cache = TTLCache(maxsize=64, ttl=120)


@dataclass
class AttachmentContext:
    """Structured file context shared across ALL teammates."""
    file_id: str
    filename: str
    type: str  # pdf | image | code | text | excel | binary | unknown
    size: int
    summary: str  # concise description for prompt injection
    chunks: list[str] = field(default_factory=list)  # text chunks (if applicable)
    extracted_entities: list[str] = field(default_factory=list)  # key terms / entities
    metadata: dict = field(default_factory=dict)

    def to_prompt_text(self, max_chars: int = 2000) -> str:
        """Convert to compact text block for LLM prompt injection."""
        parts = [f"[{self.type.upper()}: {self.filename}]"]

        if self.summary:
            parts.append(f"Summary: {self.summary}")

        if self.extracted_entities:
            entities_str = ", ".join(self.extracted_entities[:15])
            parts.append(f"Key topics: {entities_str}")

        if self.chunks:
            # Join chunks up to max_chars
            text_content = "\n".join(self.chunks)
            if len(text_content) > max_chars:
                text_content = text_content[:max_chars] + "...(truncated)"
            parts.append(f"Content:\n{text_content}")

        return "\n".join(parts)

    def to_dict(self) -> dict:
        return asdict(self)

    @staticmethod
    def from_dict(d: dict) -> "AttachmentContext":
        return AttachmentContext(**d)


def _detect_file_type(filename: str, content_bytes: bytes) -> str:
    """Detect file type from extension and magic bytes."""
    ext = os.path.splitext(filename or "")[1].lower()

    # Image
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}:
        return "image"
    # Office documents
    if ext in {".docx"}:
        return "document"
    if ext in {".pptx"}:
        return "presentation"
    # PDF
    if ext == ".pdf" or content_bytes[:4] == b"%PDF":
        return "pdf"
    # Code / text
    code_exts = {
        ".py", ".js", ".ts", ".jsx", ".tsx", ".go", ".rs", ".java", ".kt",
        ".c", ".cpp", ".h", ".hpp", ".swift", ".rb", ".php", ".sh", ".sql",
        ".vue", ".svelte",
    }
    if ext in code_exts:
        return "code"
    # Excel
    if ext in {".xls", ".xlsx", ".csv"}:
        return "excel"
    # Text
    text_exts = {
        ".txt", ".md", ".json", ".yaml", ".yml", ".html", ".htm", ".css",
        ".xml", ".toml", ".ini", ".cfg", ".conf", ".log", ".rst", ".env",
    }
    if ext in text_exts:
        return "text"
    # Try text detection
    try:
        content_bytes[:1024].decode("utf-8")
        return "text"
    except UnicodeDecodeError:
        pass
    return "binary"


def _extract_text_chunks(
    content_bytes: Optional[bytes] = None,
    max_total: int = 8000,
    text: Optional[str] = None,
) -> list[str]:
    """Split text content into manageable chunks for LLM context."""
    if text is None:
        if not content_bytes:
            return ["[Empty content]"]
        try:
            text = content_bytes.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = content_bytes.decode("latin-1")
            except Exception as e:
                logger.warning(f"Attachment service decode error: {e}")
                return ["[Binary content - not readable]"]

    if len(text) > max_total:
        text = text[:max_total]

    # Split by lines into chunks of ~50 lines
    lines = text.split("\n")
    chunks = []
    current_chunk = []
    current_len = 0

    for line in lines:
        if current_len + len(line) + 1 > 1000 or len(current_chunk) >= 50:
            if current_chunk:
                chunks.append("\n".join(current_chunk))
            current_chunk = [line]
            current_len = len(line)
        else:
            current_chunk.append(line)
            current_len += len(line) + 1

    if current_chunk:
        chunks.append("\n".join(current_chunk))

    return chunks[:10]  # Max 10 chunks


def _extract_entities(text: str) -> list[str]:
    """Extract key entities / terms from text content."""
    entities = set()

    # Capitalized words (likely proper nouns / class names)
    cap_words = re.findall(r"\b[A-Z][a-zA-Z_]{2,}\b", text)
    entities.update(cap_words[:10])

    # Chinese terms (2-8 chars)
    cn_terms = re.findall(r"[\u4e00-\u9fff]{2,8}", text)
    entities.update(cn_terms[:8])

    # Function/method names
    func_names = re.findall(r"\b(def |class |function |func )(\w+)", text)
    entities.update(name for _, name in func_names[:5])

    # File paths / URLs
    paths = re.findall(r"(?:/[a-zA-Z0-9_.-]+){2,}", text)
    entities.update(paths[:5])

    return list(entities)[:15]


def _extract_office_text(content_bytes: bytes, ext: str) -> str:
    """Extract text from office documents using specialized parsers."""
    import io
    text = ""
    try:
        if ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            text = "\n".join(texts)
        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content_bytes))
            texts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
            text = "\n".join(texts)
    except Exception as e:
        logger.warning(f"Failed to parse office document ({ext}): {e}")
    return text


def _generate_summary(filename: str, file_type: str, chunks: list[str], size: int) -> str:
    """Generate a concise summary of the file."""
    size_str = f"{size}B" if size < 1024 else f"{size/1024:.1f}KB" if size < 1024*1024 else f"{size/(1024*1024):.1f}MB"

    if file_type == "image":
        return f"Image file ({size_str})"
    if file_type == "pdf":
        lines = sum(c.count("\n") + 1 for c in chunks) if chunks else 0
        return f"PDF document ({lines} lines, {size_str})"
    if file_type == "code":
        lines = sum(c.count("\n") + 1 for c in chunks) if chunks else 0
        return f"Code file ({lines} lines, {size_str})"
    if file_type == "excel":
        return f"Spreadsheet ({size_str})"
    if file_type == "document":
        lines = sum(c.count("\n") + 1 for c in chunks) if chunks else 0
        return f"Word document ({lines} lines, {size_str})"
    if file_type == "presentation":
        return f"PowerPoint presentation ({len(chunks)} slides, {size_str})"
    if file_type == "text":
        return f"Text file ({size_str})"
    if file_type == "binary":
        return f"Binary file ({size_str}, content may not be readable)"
    return f"File: {filename} ({size_str})"


def parse_file_context(
    file_id: str,
    filename: str,
    content_bytes: bytes,
) -> AttachmentContext:
    """
    Parse file bytes into a structured AttachmentContext.
    This is the SINGLE parsing point — result is cached and shared across ALL teammates.
    """
    # Check cache first
    cached = attachment_cache.get(f"ctx:{file_id}")
    if cached is not None:
        logger.debug(f"AttachmentContext cache hit for {file_id}")
        return AttachmentContext.from_dict(cached)

    file_type = _detect_file_type(filename, content_bytes)
    chunks = []
    entities = []

    if file_type in ("text", "code"):
        chunks = _extract_text_chunks(content_bytes)
        text_for_entities = "\n".join(chunks)
        entities = _extract_entities(text_for_entities)
    elif file_type in ("excel",):
        chunks = _extract_text_chunks(content_bytes)
        text_for_entities = "\n".join(chunks)
        entities = _extract_entities(text_for_entities)
    elif file_type in ("pdf", "document", "presentation"):
        ext = os.path.splitext(filename or "")[1].lower()
        extracted = _extract_office_text(content_bytes, ext)
        if not extracted:
            extracted = "[Content could not be extracted]"
        chunks = _extract_text_chunks(text=extracted)
        text_for_entities = "\n".join(chunks)
        entities = _extract_entities(text_for_entities)
    elif file_type == "binary":
        chunks = _extract_text_chunks(content_bytes[:4096])  # Try first 4KB
        if not chunks:
            chunks = ["[Binary content]"]

    summary = _generate_summary(filename, file_type, chunks, len(content_bytes))

    ctx = AttachmentContext(
        file_id=file_id,
        filename=filename,
        type=file_type,
        size=len(content_bytes),
        summary=summary,
        chunks=chunks,
        extracted_entities=entities,
        metadata={
            "detected_type": file_type,
            "chunks_count": len(chunks),
        },
    )

    # Cache
    attachment_cache.set(f"ctx:{file_id}", ctx.to_dict())
    logger.info(f"Parsed attachment context: {filename} ({file_type}, {len(chunks)} chunks)")
    return ctx


def get_attachment_context(file_id: str) -> Optional[AttachmentContext]:
    """Retrieve cached attachment context by file_id."""
    cached = attachment_cache.get(f"ctx:{file_id}")
    if cached:
        return AttachmentContext.from_dict(cached)
    return None


def build_shared_attachment_prompt(
    attachments: list[dict],
    max_chars_per_file: int = 2000,
) -> str:
    """
    Build a combined prompt section from multiple attachments.
    Used by team_collaboration to inject into ALL teammate prompts uniformly.
    """
    if not attachments:
        return ""

    contexts = []
    for att in attachments:
        file_id = att.get("file_id") or att.get("saved_as") or att.get("filename", "")
        # Try cache first
        ctx = get_attachment_context(file_id)
        if ctx:
            contexts.append(ctx.to_prompt_text(max_chars=max_chars_per_file))
        elif att.get("llm_content"):
            # Fallback: use raw llm_content if no parsed context yet
            raw = att["llm_content"]
            if isinstance(raw, str):
                truncated = raw[:max_chars_per_file]
                if len(raw) > max_chars_per_file:
                    truncated += "...(truncated)"
                contexts.append(f"[File: {att.get('filename', '?')}]\n{truncated}")

    if not contexts:
        return ""

    return "## Shared Files Context\nThe following files were attached to the conversation. ALL teammates should reference this same context:\n\n" + "\n\n---\n\n".join(contexts)
