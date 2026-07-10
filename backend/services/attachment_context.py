"""
attachment_context.py — Shared file context pipeline for multi-teammate collaboration.

Pipeline:
  File bytes → Parse + Hash → AttachmentContext (dataclass) → DB cache → Inject into ALL teammates

All teammates receive the SAME parsed context snapshot (no individual parsing).
"""
import json
import logging
import re
import hashlib
import os
from typing import Any, Optional
from dataclasses import dataclass, field, asdict
from datetime import datetime

logger = logging.getLogger("attachment_context")


# ── AttachmentContext dataclass ──

@dataclass
class AttachmentContext:
    """Structured file context shared across ALL teammates."""
    file_id: str
    filename: str
    type: str  # pdf | image | code | text | excel | binary | unknown
    size: int
    summary: str
    chunks: list = field(default_factory=list)
    extracted_entities: list = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    version_key: str = ""  # "{file_id}:{content_hash}" — set after init

    def to_dict(self) -> dict:
        return asdict(self)

    @staticmethod
    def from_dict(d: dict) -> "AttachmentContext":
        # Backward compat: handle old keys
        if "chunks" in d and isinstance(d["chunks"], str):
            d["chunks"] = json.loads(d["chunks"])
        if "extracted_entities" in d and isinstance(d["extracted_entities"], str):
            d["extracted_entities"] = json.loads(d["extracted_entities"])
        if "metadata" in d and isinstance(d["metadata"], str):
            d["metadata"] = json.loads(d["metadata"])
        return AttachmentContext(**d)


# ── Content hashing ──

def compute_content_hash(content_bytes: bytes) -> str:
    """Compute a stable SHA-256 hash of file content."""
    return hashlib.sha256(content_bytes).hexdigest()[:16]  # 16 chars sufficient for dedup


# ── File type detection ──

def _detect_file_type(filename: str, content_bytes: bytes) -> str:
    """Detect file type from extension and magic bytes."""
    ext = os.path.splitext(filename or "")[1].lower()

    # Image
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}:
        return "image"
    # Office documents
    if ext in {".docx"}:
        return "document"
    if ext in {".pptx"}:
        return "presentation"
    # PDF
    if ext == ".pdf" or content_bytes[:4] == b"%PDF":
        return "pdf"
    # Code
    code_exts = {
        ".py", ".js", ".ts", ".jsx", ".tsx", ".go", ".rs", ".java", ".kt",
        ".c", ".cpp", ".h", ".hpp", ".swift", ".rb", ".php", ".sh", ".sql",
        ".vue", ".svelte",
    }
    if ext in code_exts:
        return "code"
    # Excel
    if ext in {".xls", ".xlsx", ".csv"}:
        return "excel"
    # Text
    text_exts = {
        ".txt", ".md", ".json", ".yaml", ".yml", ".html", ".htm", ".css",
        ".xml", ".toml", ".ini", ".cfg", ".conf", ".log", ".rst", ".env",
    }
    if ext in text_exts:
        return "text"
    # Try text detection via decode
    try:
        content_bytes[:1024].decode("utf-8")
        return "text"
    except UnicodeDecodeError:
        pass
    return "binary"


# ── Text chunking ──

def _extract_text_chunks(
    content_bytes: Optional[bytes] = None,
    max_total: int = 8000,
    text: Optional[str] = None,
) -> list:
    """Split text content into manageable chunks for LLM context."""
    if text is None:
        if not content_bytes:
            return ["[Empty content]"]
        try:
            text = content_bytes.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = content_bytes.decode("latin-1")
            except Exception as e:
                logger.warning(f"Attachment context decode error: {e}")
                return ["[Binary content - not readable]"]

    if len(text) > max_total:
        text = text[:max_total]

    lines = text.split("\n")
    chunks = []
    current_chunk = []
    current_len = 0

    for line in lines:
        if current_len + len(line) + 1 > 1000 or len(current_chunk) >= 50:
            if current_chunk:
                chunks.append("\n".join(current_chunk))
            current_chunk = [line]
            current_len = len(line)
        else:
            current_chunk.append(line)
            current_len += len(line) + 1

    if current_chunk:
        chunks.append("\n".join(current_chunk))

    return chunks[:10]  # Max 10 chunks


# ── Entity extraction ──

def _extract_entities(text: str) -> list:
    """Extract key entities / terms from text content."""
    entities = set()

    # Capitalized words (likely proper nouns / class names / acronyms)
    cap_words = re.findall(r"\b[A-Z][a-zA-Z_]{2,}\b", text)
    entities.update(cap_words[:10])

    # Chinese terms (2-8 chars)
    cn_terms = re.findall(r"[\u4e00-\u9fff]{2,8}", text)
    entities.update(cn_terms[:8])

    # Function/method names
    func_names = re.findall(r"\b(def |class |function |func )(\w+)", text)
    entities.update(name for _, name in func_names[:5])

    # File paths
    paths = re.findall(r"(?:/[a-zA-Z0-9_.-]+){2,}", text)
    entities.update(paths[:5])

    return list(entities)[:20]


def _extract_office_text(content_bytes: bytes, ext: str) -> str:
    """Extract text from office documents using specialized parsers."""
    import io
    text = ""
    try:
        if ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            text = "\n".join(texts)
        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content_bytes))
            texts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
            text = "\n".join(texts)
    except Exception as e:
        logger.warning(f"Failed to parse office document ({ext}): {e}")
    return text


# ── Summary generation ──

def _generate_summary(filename: str, file_type: str, chunks: list, size: int) -> str:
    """Generate a concise summary of the file."""
    size_str = (
        f"{size}B" if size < 1024
        else f"{size/1024:.1f}KB" if size < 1024*1024
        else f"{size/(1024*1024):.1f}MB"
    )

    if file_type == "image":
        return f"Image file ({size_str})"
    if file_type == "pdf":
        lines = sum(c.count("\n") + 1 for c in chunks) if chunks else 0
        return f"PDF document ({lines} lines, {size_str})"
    if file_type == "code":
        lines = sum(c.count("\n") + 1 for c in chunks) if chunks else 0
        return f"Code file ({lines} lines, {size_str})"
    if file_type == "excel":
        return f"Spreadsheet ({size_str})"
    if file_type == "document":
        lines = sum(c.count("\n") + 1 for c in chunks) if chunks else 0
        return f"Word document ({lines} lines, {size_str})"
    if file_type == "presentation":
        return f"PowerPoint presentation ({len(chunks)} slides, {size_str})"
    if file_type == "text":
        return f"Text file ({size_str})"
    if file_type == "binary":
        return f"Binary file ({size_str})"
    return f"File: {filename} ({size_str})"


# ── Main parsing function (SINGLE parsing point → result shared across ALL teammates) ──

def parse_file_context(
    file_id: str,
    filename: str,
    content_bytes: bytes,
) -> AttachmentContext:
    """
    Parse file bytes into a structured AttachmentContext.

    This is the SINGLE parsing point. Result is cached in DB,
    then shared across ALL teammates identically.
    """
    content_hash = compute_content_hash(content_bytes)
    version_key = f"{file_id}:{content_hash}"

    file_type = _detect_file_type(filename, content_bytes)
    chunks = []
    entities = []

    if file_type in ("text", "code"):
        chunks = _extract_text_chunks(content_bytes)
        text_for_entities = "\n".join(chunks)
        entities = _extract_entities(text_for_entities)
    elif file_type in ("excel",):
        chunks = _extract_text_chunks(content_bytes)
        text_for_entities = "\n".join(chunks)
        entities = _extract_entities(text_for_entities)
    elif file_type in ("pdf", "document", "presentation"):
        ext = os.path.splitext(filename or "")[1].lower()
        extracted = _extract_office_text(content_bytes, ext)
        if not extracted:
            extracted = "[Content could not be extracted]"
        chunks = _extract_text_chunks(text=extracted)
        text_for_entities = "\n".join(chunks) if chunks else ""
        entities = _extract_entities(text_for_entities)
    elif file_type == "binary":
        try:
            sample = content_bytes[:4096].decode("utf-8", errors="ignore")
            if any(c.isprintable() for c in sample[:100]):
                chunks = _extract_text_chunks(content_bytes[:4096])
                entities = _extract_entities("\n".join(chunks))
        except Exception as e:
            logger.debug(f"Binary content entity extraction skipped: {e}")
        if not chunks:
            chunks = ["[Binary content]"]

    summary = _generate_summary(filename, file_type, chunks, len(content_bytes))

    ctx = AttachmentContext(
        file_id=file_id,
        filename=filename,
        type=file_type,
        size=len(content_bytes),
        summary=summary,
        chunks=chunks,
        extracted_entities=entities,
        metadata={
            "filename": filename,
            "detected_type": file_type,
            "chunks_count": len(chunks),
        },
        version_key=version_key,
    )

    logger.info(f"Parsed attachment context: {filename} ({file_type}, {len(chunks)} chunks, {version_key})")
    return ctx


# ── DB persistence ──

async def _db_save_context(ctx: AttachmentContext) -> None:
    """Save AttachmentContext to DB (INSERT OR IGNORE on version_key)."""
    try:
        from backend.database import async_session
        from backend.models import AttachmentContextModel

        async with async_session() as sess:
            row = AttachmentContextModel(
                file_id=ctx.file_id,
                content_hash=ctx.version_key.split(":")[1] if ":" in ctx.version_key else "",
                type=ctx.type,
                summary=ctx.summary,
                chunks_json=json.dumps(ctx.chunks),
                entities_json=json.dumps(ctx.extracted_entities),
                metadata_json=ctx.metadata,
            )
            sess.add(row)
            await sess.commit()
            logger.info(f"AttachmentContext saved to DB: {ctx.version_key}")
    except Exception as e:
        logger.warning(f"Failed to save AttachmentContext to DB (non-fatal): {e}")


async def _db_find_context_by_file_id(file_id: str, content_hash: str) -> Optional[AttachmentContext]:
    """Look up AttachmentContext from DB by file_id and content hash."""
    try:
        from backend.database import async_session
        from backend.models import AttachmentContextModel
        from sqlalchemy import select

        async with async_session() as sess:
            result = await sess.execute(
                select(AttachmentContextModel).where(
                    AttachmentContextModel.file_id == file_id,
                    AttachmentContextModel.content_hash == content_hash,
                )
            )
            row = result.scalar_one_or_none()
            if row:
                return AttachmentContext(
                    file_id=row.file_id,
                    filename=row.metadata_json.get("filename", "unknown") if row.metadata_json else "unknown",
                    type=row.type or "text",
                    size=int(row.metadata_json.get("size", 0)) if row.metadata_json else 0,
                    summary=row.summary or "",
                    chunks=json.loads(row.chunks_json) if row.chunks_json else [],
                    extracted_entities=json.loads(row.entities_json) if row.entities_json else [],
                    metadata=json.loads(row.metadata_json) if row.metadata_json else {},
                    version_key=f"{row.file_id}:{row.content_hash}",
                )
    except Exception as e:
        logger.warning(f"DB lookup failed for {file_id}:{content_hash}: {e}")
    return None
